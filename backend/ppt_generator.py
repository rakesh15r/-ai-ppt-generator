from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def build_presentation(slides, topic_hint="Class Slides"):
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    for index, slide_data in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        _style_slide_background(slide, slide_data.get("type", "normal"))
        _render_title(slide, slide_data.get("title", f"Slide {index}"))
        _render_points(slide, slide_data.get("points", []))
        _render_footer(slide, topic_hint, index, slide_data.get("type", "normal"))

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output


def _style_slide_background(slide, slide_type):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(246, 248, 252)

    accent_color = _accent_color(slide_type)

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.55),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(15, 23, 42)
    banner.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.85),
        Inches(0.78),
        Inches(1.65),
        Inches(0.34),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()


def _render_title(slide, title):
    title_box = slide.shapes.title
    title_box.left = Inches(0.7)
    title_box.top = Inches(0.75)
    title_box.width = Inches(11.7)
    title_box.height = Inches(0.9)

    text_frame = title_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.alignment = PP_ALIGN.LEFT
    font = paragraph.runs[0].font
    font.name = "Aptos Display"
    font.size = Pt(26)
    font.bold = True
    font.color.rgb = RGBColor(15, 23, 42)


def _render_points(slide, points):
    body = slide.placeholders[1]
    body.left = Inches(0.95)
    body.top = Inches(1.8)
    body.width = Inches(11.3)
    body.height = Inches(4.95)

    text_frame = body.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    for index, point in enumerate(points):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = point
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(19)
        paragraph.font.color.rgb = RGBColor(30, 41, 59)


def _render_footer(slide, topic_hint, slide_number, slide_type):
    footer = slide.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(11.8), Inches(0.3))
    text_frame = footer.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    slide_label = slide_type.title() if slide_type else "Normal"
    paragraph.text = f"{topic_hint} | Slide {slide_number} | {slide_label}"
    paragraph.alignment = PP_ALIGN.LEFT
    font = paragraph.runs[0].font
    font.name = "Aptos"
    font.size = Pt(10)
    font.color.rgb = RGBColor(100, 116, 139)


def _accent_color(slide_type):
    if slide_type == "explanation":
        return RGBColor(245, 158, 11)
    if slide_type == "data":
        return RGBColor(34, 197, 94)
    return RGBColor(14, 165, 233)
